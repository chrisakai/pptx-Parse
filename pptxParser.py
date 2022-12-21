# -*- coding:utf-8 -*-
import json
import os.path
import shutil
import uuid
from datetime import date, datetime

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE as mst

# EMUs(English Metric Units)
# 1EMUs= 1/914400英寸
# EMUS=像素*914400/分辨率
# w像素=EMUS/914400*分辨率width
# h像素=EMUS/914400*分辨率height
# https://www.cnblogs.com/ac1985482/p/4097666.html
# screen_width = 640
# screen_height = 480

# EMU单位到像素单位转换
to_pixel = 914400 / 96

pptxParsePath = 'pptxParse'


# 日期类型转化
class ComplexEncoder(json.JSONEncoder):
    def default(self, obj):
        if isinstance(obj, datetime):
            return obj.strftime('%Y-%m-%d %H:%M:%S')
        elif isinstance(obj, date):
            return obj.strftime('%Y-%m-%d')
        else:
            return json.JSONEncoder.default(self, obj)


def toJson(info):
    # jsObj = json.dumps(info)
    # with open("pptx-parse.json", "w", encoding='utf-8') as f:
    #     f.write(jsObj)
    #     f.close()
    json.dump(info, open("pptx-parse.json", 'w', encoding='utf-8'), indent=2, ensure_ascii=False, cls=ComplexEncoder)


def mycopyfile(srcfile, dstpath):  # 复制函数
    if not os.path.isfile(srcfile):
        print("%s not exist!" % (srcfile))
    else:
        fpath, fname = os.path.split(srcfile)  # 分离文件名和路径
        if not os.path.exists(dstpath):
            os.makedirs(dstpath)  # 创建路径
        shutil.copy(srcfile, dstpath + fname)  # 复制文件
        print("copy %s -> %s" % (srcfile, dstpath + fname))


#
# 把Shape中数据，依据Shape类型转成相应Json数据
#
def toJsonObject(shape, shapeType, fileName):
    left_px = round(shape.left / to_pixel, 2)
    top_cm = round(shape.top / to_pixel, 2)
    width_cm = round(shape.width / to_pixel, 2)
    height_cm = round(shape.height / to_pixel, 2)
    jsonObj = {
        'x': left_px,
        'y': top_cm,
        'width': width_cm,
        'height': height_cm,
        'source': fileName
    }
    if shapeType == mst.PICTURE:
        jsonObj['source'] = fileName

    return jsonObj


# 解析pptx
def pptxParse(pptxName):
    slides = []
    slideCount = 1

    pptx = Presentation(pptxName)

    print("############Header#############")
    count = len(pptx.slides)
    print(count)
    pptxname = pptxName.split("/")[-1]
    print(pptxname)
    version = pptx.core_properties.revision
    print(version)
    importTime = pptx.core_properties.created
    print(importTime)
    modifyTime = pptx.core_properties.modified
    print(modifyTime)
    screen_height = pptx.slide_height / to_pixel
    print(screen_height)
    screen_width = pptx.slide_width / to_pixel
    print(screen_width)
    print("############Header#############")

    for slide in pptx.slides:
        shapes = []
        for shape in slide.shapes:
            print(shape.shape_type)
            # 箭头线
            if shape.shape_type == 9:
                # 引导线开始坐标
                # print('line {0}'.format(round(shape.begin_x / to_pixel, 2)));
                # 文本框原点坐标，宽高
                # Cx = (long)bm.Width * (long)((float)914400 / bm.HorizontalResolution);

                begin_x = round(shape.begin_x / to_pixel, 2)
                begin_y = round(shape.begin_y / to_pixel, 2)
                end_x = round(shape.end_x / to_pixel, 2)
                end_y = round(shape.end_y / to_pixel, 2)
                jsonObj = {
                    'begin_x': begin_x,
                    'begin_y': begin_y,
                    'end_x': end_x,
                    'end_y': end_y,
                    'text': text_frame.text,
                }
                if bool(shape.click_action.hyperlink.address is not None):
                    if uri != None:
                        jsonObj['uri'] = uri

                data = {'type': shape.shape_type, 'hasevent': bool(shape.click_action.hyperlink.address is not None),
                        'properties': jsonObj}
                print("++++++++++++Line++++++++++++")
                print(bool(shape.click_action.hyperlink.address is not None))
                print("++++++++++++Line++++++++++++")
                shapes.append(data)
            # 图片
            elif shape.shape_type == mst.PICTURE:
                # cv2.imshow(shape.image.blob);
                # print(shape.image.d content_type);
                guid = str(uuid.uuid4()).replace("-", "")
                # 删除'imgs/' + 'shape.image.filename + '
                if not os.path.exists('pptxParse/imgs/'):
                    os.makedirs('pptxParse/imgs/')
                fileName = 'pptxParse/imgs/' + guid + '.png'
                imgFile = open(fileName, 'wb')
                imgFile.write(shape.image.blob)
                imgFile.close()
                relativePath = 'imgs/' + guid + '.png'
                jsonObj = toJsonObject(shape, mst.PICTURE, relativePath)
                if bool(shape.click_action.hyperlink.address is not None):
                    uri = shape.click_action.hyperlink.address
                    jsonObj['uri'] = 'videos/' + uri
                data = {'type': shape.shape_type, 'hasevent': bool(shape.click_action.hyperlink.address is not None),
                        'properties': jsonObj}
                print("============Picture============")
                print(bool(shape.click_action.hyperlink.address is not None))
                print("============Picture============")
                shapes.append(data)
            # 自动形状
            elif shape.shape_type == mst.AUTO_SHAPE:
                action_setting = shape.click_action
                uri = ''
                if action_setting != None and action_setting.hyperlink.address != None:
                    uri = action_setting.hyperlink.address

                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    paragraphs = text_frame.paragraphs
                    for graph in paragraphs:
                        for run in graph.runs:
                            uri = run.hyperlink.address

                            # print(uri)

                    # 文本框原点坐标，宽高
                    # Cx = (long)bm.Width * (long)((float)914400 / bm.HorizontalResolution);
                    left_px = round(shape.left / to_pixel, 2)
                    top_cm = round(shape.top / to_pixel, 2)
                    width_cm = round(shape.width / to_pixel, 2)
                    height_cm = round(shape.height / to_pixel, 2)
                    jsonObj = {
                        'x': left_px,
                        'y': top_cm,
                        'width': width_cm,
                        'height': height_cm,
                        'text': text_frame.text,
                    }
                    if bool(shape.click_action.hyperlink.address is not None):
                        if uri != None:
                            jsonObj['uri'] = uri
                    data = {'type': shape.shape_type,
                            'hasevent': bool(shape.click_action.hyperlink.address is not None), 'properties': jsonObj}
                    print("------------AutoShape------------")
                    print(bool(shape.click_action.hyperlink.address is not None))
                    print("------------AutoShape------------")
                    shapes.append(data)
            # 文本框
            elif shape.shape_type == mst.TEXT_BOX:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    paragraphs = text_frame.paragraphs
                    uri = ''
                    paragraphs = text_frame.paragraphs
                    for graph in paragraphs:
                        for run in graph.runs:
                            uri = run.hyperlink.address

                    # 文本框原点坐标，宽高
                    # Cx = (long)bm.Width * (long)((float)914400 / bm.HorizontalResolution);
                    left_px = round(shape.left / to_pixel, 2)
                    top_cm = round(shape.top / to_pixel, 2)
                    width_cm = round(shape.width / to_pixel, 2)
                    height_cm = round(shape.height / to_pixel, 2)
                    jsonObj = {
                        'x': left_px,
                        'y': top_cm,
                        'width': width_cm,
                        'height': height_cm,
                        'text': text_frame.text,
                    }
                    if uri != None:
                        jsonObj['uri'] = 'videos/' + uri
                        data = {'type': shape.shape_type,
                                'hasevent': True,
                                'properties': jsonObj}
                    else:
                        data = {'type': shape.shape_type,
                                'hasevent': False,
                                'properties': jsonObj}
                    print("***********Text*************")
                    print(bool(shape.click_action.hyperlink.address is not None))
                    print("***********Text*************")
                    shapes.append(data)
        slidePage = {'slide': "第" + str(slideCount) + "页", 'shapes': shapes}
        slideCount = slideCount + 1
        slides.append(slidePage)

    Dict = {
        "slideCount": count,
        "filename": pptxname,
        "version": version,
        "importTime": importTime,
        "modifyTime": modifyTime,
        "width": screen_width,
        "height": screen_height,
        "slides": slides
    }

    toJson(Dict)

    mycopyfile('pptx-parse.json', pptxParsePath + '/')
